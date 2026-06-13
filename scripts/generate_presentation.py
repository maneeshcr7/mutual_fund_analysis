"""
Bluestock Mutual Fund Analysis - Presentation Generator
Generates a 12-slide PowerPoint presentation for stakeholder presentation.

Slides:
1. Title
2. Problem & Objective
3. Data Sources
4. Architecture
5. EDA Highlights (Part 1)
6. EDA Highlights (Part 2)
7. Performance Metrics (Part 1)
8. Performance Metrics (Part 2)
9. Dashboard Screenshots (Part 1)
10. Dashboard Screenshots (Part 2)
11. Key Findings
12. Thank You
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Workspace paths
WORKSPACE = Path(__file__).parent.parent
REPORTS_DIR = WORKSPACE / "reports"
FIGURES_DIR = REPORTS_DIR / "figures"


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide to the presentation."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    
    # Set subtitle
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x5c, 0x6b, 0xc0)
    
    return slide


def add_content_slide(prs, title, bullets, notes=""):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x28, 0x35, 0x93)
    
    # Add bullets
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.space_after = Pt(12)
        p.level = 0
    
    # Add notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_chart_slide(prs, title, fig_file, notes=""):
    """Add a slide with a chart/image."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0x28, 0x35, 0x93)
    
    # Add chart image
    fig_path = FIGURES_DIR / fig_file
    if fig_path.exists():
        try:
            slide.shapes.add_picture(
                str(fig_path),
                Inches(1.0), Inches(1.5),
                width=Inches(11)
            )
        except Exception as e:
            # Add placeholder text if image fails
            placeholder = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(2))
            placeholder.text_frame.text = f"[Chart: {title}]"
    else:
        # Add placeholder text if file not found
        placeholder = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(2))
        placeholder.text_frame.text = f"[Chart: {title}]"
    
    # Add notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_two_chart_slide(prs, title, fig_file1, fig_file2, label1="", label2=""):
    """Add a slide with two charts side by side."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0x28, 0x35, 0x93)
    
    # Add first chart
    fig_path1 = FIGURES_DIR / fig_file1
    if fig_path1.exists():
        try:
            slide.shapes.add_picture(
                str(fig_path1),
                Inches(0.5), Inches(1.2),
                width=Inches(6),
                height=Inches(4.5)
            )
        except Exception:
            pass
    
    # Add second chart
    fig_path2 = FIGURES_DIR / fig_file2
    if fig_path2.exists():
        try:
            slide.shapes.add_picture(
                str(fig_path2),
                Inches(6.8), Inches(1.2),
                width=Inches(6),
                height=Inches(4.5)
            )
        except Exception:
            pass
    
    return slide


def generate_presentation():
    """Generate the 12-slide presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Bluestock Mutual Fund Analysis",
        "Capstone Project — Performance Analytics & Advanced Analysis"
    )
    
    # Slide 2: Problem & Objective
    add_content_slide(
        prs,
        "Problem & Objective",
        [
            "Problem: Investors lack data-driven tools for mutual fund selection",
            "Objective: Build comprehensive analytics platform for 40 MF schemes",
            "Scope: ETL pipeline, performance metrics, risk analysis, recommender system",
            "Deliverables: Report, Dashboard, Presentation, GitHub repository",
            "Tools: Python, SQLite, Power BI/Tableau, scikit-learn"
        ],
        "Introduce the problem statement and project objectives"
    )
    
    # Slide 3: Data Sources
    add_content_slide(
        prs,
        "Data Sources",
        [
            "10 datasets covering 87,493+ records",
            "40 mutual fund schemes across 10 fund houses",
            "Analysis period: January 2022 – May 2026",
            "NAV history: 46,000 daily records",
            "Investor transactions: 32,778 records",
            "Benchmark indices: 7 indices (Nifty 50, 100, 500, Midcap, SmallCap, Gilt, Liquid)"
        ],
        "Overview of data sources and dataset sizes"
    )
    
    # Slide 4: Architecture
    add_content_slide(
        prs,
        "Architecture",
        [
            "Extract: Load 10 CSV datasets with pathlib.Path",
            "Transform: Clean, validate, forward-fill, deduplicate",
            "Load: SQLite star schema (2 dim + 9 fact tables)",
            "Analyze: Python (pandas, numpy, scipy, scikit-learn)",
            "Visualize: matplotlib, seaborn, plotly",
            "Dashboard: Power BI / Tableau with slicers",
            "Report: Automated PDF and PPTX generation"
        ],
        "Technical architecture overview"
    )
    
    # Slide 5: EDA Highlights (Part 1)
    add_two_chart_slide(
        prs,
        "EDA Highlights — NAV Trends & AUM Growth",
        "01_nav_trend_all_schemes.png",
        "02_aum_growth_by_fund_house.png",
        "NAV Trends",
        "AUM Growth"
    )
    
    # Slide 6: EDA Highlights (Part 2)
    add_two_chart_slide(
        prs,
        "EDA Highlights — SIP Inflows & Demographics",
        "03_sip_inflow_timeseries.png",
        "05_investor_demographics.png",
        "SIP Inflows",
        "Demographics"
    )
    
    # Slide 7: Performance Metrics (Part 1)
    add_two_chart_slide(
        prs,
        "Performance Metrics — CAGR & Sharpe Ratio",
        "11_cagr_comparison.png",
        "12_sharpe_ratio_ranking.png",
        "CAGR Comparison",
        "Sharpe Ratio"
    )
    
    # Slide 8: Performance Metrics (Part 2)
    add_two_chart_slide(
        prs,
        "Performance Metrics — Alpha/Beta & Drawdown",
        "14_alpha_beta_scatter.png",
        "15_drawdown_curves.png",
        "Alpha vs Beta",
        "Drawdown Curves"
    )
    
    # Slide 9: Dashboard Screenshots (Part 1)
    add_two_chart_slide(
        prs,
        "Dashboard — Overview & Fund Performance",
        "16_fund_scorecard.png",
        "17_benchmark_comparison_3yr.png",
        "Fund Scorecard",
        "Benchmark Comparison"
    )
    
    # Slide 10: Dashboard Screenshots (Part 2)
    add_two_chart_slide(
        prs,
        "Dashboard — Risk Analysis & Cohort",
        "20_var_cvar_comparison.png",
        "22_cohort_analysis.png",
        "VaR Analysis",
        "Cohort Analysis"
    )
    
    # Slide 11: Key Findings
    add_content_slide(
        prs,
        "Key Findings",
        [
            "Top funds: Mirae Asset Large Cap (86.25), ICICI Pru Midcap (82.25), Kotak Flexicap (82.00)",
            "Mid-cap funds delivered highest risk-adjusted returns (avg Sharpe 1.1)",
            "26-35 age group drives 42% of investment volume",
            "SIP retention strong: avg 8+ transactions over 200+ days",
            "Direct plans outperform regular by 0.5-1.0% (lower expense ratio)",
            "VaR ranges from 0.1% (liquid) to 3.5% (small cap) at 95% confidence"
        ],
        "Summary of key findings from the analysis"
    )
    
    # Slide 12: Thank You / Q&A
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add thank you text
    thank_you_box = slide.shapes.add_textbox(Inches(3), Inches(2.5), Inches(7), Inches(2))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You!"
    thank_you_frame.paragraphs[0].font.size = Pt(54)
    thank_you_frame.paragraphs[0].font.bold = True
    thank_you_frame.paragraphs[0].font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    thank_you_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(3), Inches(4.5), Inches(7), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions & Discussion"
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(0x5c, 0x6b, 0xc0)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add contact info
    contact_box = slide.shapes.add_textbox(Inches(3), Inches(5.8), Inches(7), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "GitHub: github.com/maneeshreddy/bluestock-mf-capstone"
    contact_frame.paragraphs[0].font.size = Pt(16)
    contact_frame.paragraphs[0].font.color.rgb = RGBColor(0x61, 0x61, 0x61)
    contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = REPORTS_DIR / "Bluestock_MF_Presentation.pptx"
    prs.save(str(output_path))
    print(f"Bluestock_MF_Presentation.pptx generated successfully!")
    print(f"Saved to: {output_path}")
    
    return True


if __name__ == "__main__":
    os.chdir(WORKSPACE)
    generate_presentation()
